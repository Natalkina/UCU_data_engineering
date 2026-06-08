from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)
        p2.alignment = PP_ALIGN.CENTER


def add_content_slide(title, bullets, font_size=18):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.5), Inches(5.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)


def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    cols_count = len(headers)
    rows_count = len(rows) + 1
    col_width = 11.5 / cols_count

    table_shape = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.1), Inches(11.5), Inches(5.5))
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY


# === SLIDES ===

# 1. Title
add_title_slide(
    "Direct Marketing Materials\nCollection & Annotation Pipeline",
    "Solution Design"
)

# 2. Task 1 - Scale & Volume
add_table_slide(
    "Task 1: Assumptions on Scale & Volume",
    ["Parameter", "Assumption"],
    [
        ["Daily ingestion volume", "~5,000 marketing materials"],
        ["Average file size", "~2 MB (scanned prints + digital)"],
        ["Daily storage growth", "~10 GB/day raw, ~15 GB/day with derivatives"],
        ["Monthly volume", "~150K items, ~450 GB"],
        ["Regions covered", "50 states, ~200 metro areas"],
        ["Industries tracked", "20+ (retail, real estate, food, auto, healthcare...)"],
        ["Annotation turnaround", "< 24 hours from ingestion"],
        ["Concurrent annotators", "10–30 human annotators"],
        ["Query latency", "< 2 seconds for metadata queries"],
        ["Retention", "2 years minimum"],
    ]
)

# 3. Functional Requirements
add_table_slide(
    "Functional Requirements",
    ["Component", "Description"],
    [
        ["Ingestion", "Images/PDFs via mobile app, email, web scraping, partner APIs. Batch + real-time."],
        ["Pre-processing", "Normalize formats, resize, de-duplicate (perceptual hashing), validate."],
        ["OCR", "Extract all text: headlines, body copy, contact info, offers, disclaimers."],
        ["Metadata Extraction", "Classify industry, detect logos/brands, colors, layout, geography, language."],
        ["Annotation", "Human-in-the-loop: creative type, tone, audience, CTA style, effectiveness."],
        ["Storage", "Raw files, processed files, OCR text, ML metadata, annotations — queryable."],
        ["Search & Analytics", "Full-text search on OCR, filter by industry/region/date, dashboards."],
    ]
)

# 4. Non-Functional Requirements
add_table_slide(
    "Non-Functional Requirements",
    ["Requirement", "Target"],
    [
        ["Throughput", "5K items/day sustained; burst to 20K during campaigns"],
        ["Availability", "99.5% ingestion/UI; 99.9% storage"],
        ["Scalability", "Horizontal scaling of OCR/ML workers"],
        ["Security", "Encryption at rest/transit; RBAC; audit logs; GDPR PII redaction"],
        ["Cost", "< $3,000/month at steady-state"],
        ["Latency", "Ingestion-to-searchable < 15 min (real-time); < 4 hrs (batch)"],
    ]
)

# 5. Architecture Diagram (text-based)
add_content_slide(
    "Task 2: Raw Solution Architecture",
    [
        "Data Sources → Ingestion Gateway → Raw Object Store",
        "",
        "Orchestrator triggers parallel processing:",
        "  • Pre-processing (normalize, resize, dedup, validate)",
        "  • OCR Engine (text extraction, layout detection)",
        "  • ML/NLP Metadata (industry classify, logo detect)",
        "  • Annotation Tool (human-in-the-loop, QA review)",
        "",
        "Results → Processed Object Store → Metadata Database",
        "",
        "Search & Analytics Layer (full-text index + dashboards + API)",
    ],
    font_size=18
)

# 6. Data Flow
add_content_slide(
    "Data Flow Summary",
    [
        "1. Sources → Ingestion Gateway → Raw Store",
        "2. Orchestrator triggers: Pre-process → OCR → ML Metadata (parallel)",
        "3. Results land in Metadata DB",
        "4. Annotation Tool pulls items needing human review → writes back",
        "5. Search/Analytics layer indexes Metadata DB for end-user queries",
    ],
    font_size=20
)

# 7. Service Analysis - OCR
add_table_slide(
    "Task 3: OCR Engine — Options",
    ["Option", "Pros", "Cons"],
    [
        ["A: Tesseract + EasyOCR\n(Open-Source)", "Free, 100+ languages,\nno lock-in", "High ops burden,\nGPU infra needed"],
        ["B: Google Cloud Vision API\n(SaaS)", "Production-grade, auto-scales,\nlogo+label detection", "Medium lock-in,\nAPI-specific format"],
        ["C: Build from Scratch", "Full control", "3-6 months, $60K-$120K,\n+ GPU costs"],
    ]
)

# 8. OCR Recommendation
add_content_slide(
    "OCR Recommendation: Google Cloud Vision API",
    [
        "✓ Best accuracy on stylized marketing text",
        "✓ Zero operational burden — fully managed",
        "✓ Cost-effective: ~$225/month at 150K images ($1.50/1000 images)",
        "✓ Includes logo detection, label detection, layout analysis",
        "✓ Auto-scales for burst periods",
    ],
    font_size=20
)

# 9. Annotation Tool
add_table_slide(
    "Annotation Tool — Options",
    ["Option", "Pros", "Cons"],
    [
        ["A: Label Studio\n(Self-hosted on GKE)", "Free, customizable,\nML pre-labeling, REST API", "Medium ops burden,\nself-managed infra"],
        ["B: Vertex AI Data Labeling\n(SaaS)", "Managed workforce,\nbuilt-in QA, IAM", "High cost ($0.08-$0.12/label),\nhigh lock-in"],
        ["C: Build from Scratch", "Fully custom UI", "2-4 months, $40K-$80K"],
    ]
)

# 10. Annotation Recommendation
add_content_slide(
    "Annotation Recommendation: Label Studio on GKE",
    [
        "✓ Cost: $150/month vs $12,000/month (Vertex AI) vs $5,000/month (Labelbox)",
        "✓ Fully customizable taxonomy for marketing-specific labels",
        "✓ Supports custom ML backends for pre-annotation (Vision API results)",
        "✓ Zero vendor lock-in — data in standard formats (JSON, COCO)",
        "✓ GKE handles autoscaling for burst annotation periods",
    ],
    font_size=20
)

# 11. Orchestration
add_content_slide(
    "Workflow Orchestration: Cloud Composer (Managed Airflow)",
    [
        "Options considered:",
        "  A: Apache Airflow (self-hosted) — medium ops burden",
        "  B: Cloud Composer (managed) — low ops, native GCP integration",
        "  C: Custom scheduler — almost never justified",
        "",
        "Recommendation: Cloud Composer",
        "  ✓ Zero-ops orchestration",
        "  ✓ Native GCP operators (GCS, Vision API, BigQuery)",
        "  ✓ DAGs remain portable (standard Airflow)",
        "  ✓ Auto-scaling workers",
    ],
    font_size=18
)

# 12. Task 4 - Detailed Comparison
add_table_slide(
    "Task 4: Annotation Tool — Detailed Comparison",
    ["Criteria", "Label Studio (GKE)", "Vertex AI", "Labelbox"],
    [
        ["Monthly Cost", "~$150", "~$12,000", "~$5,000"],
        ["Pricing Model", "Free + infra", "$0.08–$0.12/label", "$500/user/month"],
        ["Custom Taxonomy", "✅ Fully flexible", "⚠️ Limited", "✅ Flexible"],
        ["ML Pre-labeling", "✅ Custom backend", "✅ AutoML", "✅ Model-assisted"],
        ["QA Workflow", "✅ Basic", "✅ Built-in", "✅ Advanced"],
        ["Lock-in Risk", "Low", "Medium-High", "Medium"],
        ["Data Residency", "You control", "GCP regions", "Labelbox cloud"],
        ["API Access", "REST API", "gRPC/REST", "GraphQL"],
    ]
)

# 13. Final Architecture
add_content_slide(
    "Task 5: Final Architecture — GCP Services",
    [
        "• Event Bus: Cloud Pub/Sub — handles 20K burst ingestion",
        "• Object Storage: GCS — $0.02/GB, 99.999% durability",
        "• Orchestration: Cloud Composer — managed Airflow",
        "• Pre-processing: Cloud Functions — serverless, scales to 0",
        "• OCR: Cloud Vision API — best accuracy, $1.50/1K images",
        "• ML Classification: Vertex AI — custom industry classifier",
        "• Annotation: Label Studio on GKE — $150 vs $12K/mo",
        "• Metadata DB: Cloud SQL (PostgreSQL) — ACID, managed HA",
        "• Search: Elasticsearch on GKE — full-text OCR search <2s",
        "• Dashboards: Looker Studio — free, native GCP",
        "• Search API: Cloud Run — scales to 0",
    ],
    font_size=16
)

# 14. Cost Summary
add_table_slide(
    "Estimated Monthly Cost",
    ["Service", "Monthly Cost"],
    [
        ["GCS (500 GB)", "$10"],
        ["Cloud Pub/Sub", "$5"],
        ["Cloud Composer (small)", "$400"],
        ["Cloud Functions", "$20"],
        ["Cloud Vision API (150K calls)", "$225"],
        ["Vertex AI endpoint", "$200"],
        ["Cloud SQL (PostgreSQL)", "$50"],
        ["GKE cluster (Label Studio + ES)", "$250"],
        ["Cloud Run", "$15"],
        ["Looker Studio", "$0"],
        ["TOTAL", "~$1,175/month ✅ (target < $3,000)"],
    ]
)

# 15. Summary
add_title_slide(
    "Summary",
    "Total: ~$1,175/month — well within $3,000 budget\n"
    "Key choices: Cloud Vision API • Label Studio • Cloud Composer • GKE"
)

# Save
output_path = "/Users/nataliasokil/PycharmProjects/UCU_data_engineering/automated_data_pipelines/marketing_materials_pipeline/solution_design.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
